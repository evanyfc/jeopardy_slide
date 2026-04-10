#!/usr/bin/env python3
"""
Jeopardy PowerPoint Generator

Creates a Jeopardy-style PowerPoint presentation (.pptx) with:
  - A main board slide: 6 categories x 5 dollar-value cells
  - 30 Clue slides (one per cell)
  - 30 Answer slides (one per clue)
  - Hyperlink navigation: Board -> Clue -> Answer -> Board

Usage:
    python jeopardy.py [--categories CAT1 ... CAT6] [--round single|double|triple]
                       [--values V1 V2 V3 V4 V5] [--output FILE]
"""

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
JEOPARDY_BLUE = RGBColor(0x06, 0x0C, 0xE9)
DARK_BLUE = RGBColor(0x00, 0x00, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xFF, 0xCC, 0x00)

# ---------------------------------------------------------------------------
# Dollar values per round type
# ---------------------------------------------------------------------------
ROUND_VALUES = {
    "single": [200, 400, 600, 800, 1000],
    "double": [400, 800, 1200, 1600, 2000],
    "triple": [600, 1200, 1800, 2400, 3000],
}

DEFAULT_CATEGORIES = [
    "Category 1",
    "Category 2",
    "Category 3",
    "Category 4",
    "Category 5",
    "Category 6",
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_cell_bg(shape, color):
    """Fill a shape with a solid colour."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_border(shape, color=WHITE, width_pt=1):
    """Set a shape's border colour and width."""
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)


def add_cell(slide, text, left, top, width, height,
             font_size=24, font_color=WHITE, bg_color=JEOPARDY_BLUE,
             bold=False, border_color=WHITE):
    """
    Add a rectangle shape with centred text to *slide* and return the shape.

    The shape is suitable for click-action hyperlinks.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    _set_cell_bg(shape, bg_color)
    _set_border(shape, border_color)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Remove default margins so text fills the cell properly
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER

    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold

    return shape


def set_slide_hyperlink(shape, target_slide):
    """Make *shape* a hyperlink that jumps to *target_slide*."""
    shape.click_action.target_slide = target_slide


def set_bg(slide, color):
    """Set the background fill of *slide* to *color*."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_board_slide(prs, categories, values):
    """
    Create the main Jeopardy board slide.

    Returns:
        (slide, value_shapes)  where *value_shapes* is a list of lists:
        value_shapes[row_idx][col_idx] -> shape for that cell.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_bg(slide, JEOPARDY_BLUE)

    sw = prs.slide_width
    sh = prs.slide_height
    num_cols = len(categories)
    num_rows = len(values)

    col_w = sw // num_cols
    header_h = int(sh * 0.18)
    row_h = (sh - header_h) // num_rows

    # Category header row
    for col, cat in enumerate(categories):
        add_cell(
            slide, cat,
            left=col * col_w, top=0, width=col_w, height=header_h,
            font_size=18, font_color=WHITE, bg_color=JEOPARDY_BLUE,
            bold=True,
        )

    # Dollar-value rows
    value_shapes = []
    for row, val in enumerate(values):
        row_shapes = []
        for col in range(num_cols):
            shape = add_cell(
                slide, f"${val:,}",
                left=col * col_w,
                top=header_h + row * row_h,
                width=col_w,
                height=row_h,
                font_size=32, font_color=GOLD, bg_color=JEOPARDY_BLUE,
                bold=True,
            )
            row_shapes.append(shape)
        value_shapes.append(row_shapes)

    return slide, value_shapes


def build_clue_slide(prs, category, value):
    """
    Create a Clue slide.

    Returns:
        (slide, answer_button_shape)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, JEOPARDY_BLUE)

    sw = prs.slide_width
    sh = prs.slide_height

    header_h = int(sh * 0.15)
    btn_h = int(sh * 0.12)
    clue_h = sh - header_h - btn_h

    # Header: category – dollar value
    add_cell(
        slide, f"{category}  •  ${value:,}",
        left=0, top=0, width=sw, height=header_h,
        font_size=26, font_color=GOLD, bg_color=JEOPARDY_BLUE, bold=True,
        border_color=GOLD,
    )

    # Clue body (blank — user fills in)
    add_cell(
        slide, "[ Enter Clue Here ]",
        left=0, top=header_h, width=sw, height=clue_h,
        font_size=36, font_color=WHITE, bg_color=JEOPARDY_BLUE,
        border_color=WHITE,
    )

    # "Show Answer" button
    btn_w = int(sw * 0.40)
    btn_left = (sw - btn_w) // 2
    answer_btn = add_cell(
        slide, "Show Answer  →",
        left=btn_left, top=header_h + clue_h, width=btn_w, height=btn_h,
        font_size=20, font_color=WHITE, bg_color=DARK_BLUE, bold=True,
        border_color=WHITE,
    )

    return slide, answer_btn


def build_answer_slide(prs, category, value):
    """
    Create an Answer slide.

    Returns:
        (slide, back_to_board_button_shape)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, JEOPARDY_BLUE)

    sw = prs.slide_width
    sh = prs.slide_height

    header_h = int(sh * 0.15)
    btn_h = int(sh * 0.12)
    answer_h = sh - header_h - btn_h

    # Header
    add_cell(
        slide, f"Answer  •  {category}  •  ${value:,}",
        left=0, top=0, width=sw, height=header_h,
        font_size=22, font_color=GOLD, bg_color=JEOPARDY_BLUE, bold=True,
        border_color=GOLD,
    )

    # Answer body (blank — user fills in)
    add_cell(
        slide, "[ Enter Answer Here ]",
        left=0, top=header_h, width=sw, height=answer_h,
        font_size=36, font_color=WHITE, bg_color=JEOPARDY_BLUE,
        border_color=WHITE,
    )

    # "Back to Board" button
    btn_w = int(sw * 0.40)
    btn_left = (sw - btn_w) // 2
    board_btn = add_cell(
        slide, "←  Back to Board",
        left=btn_left, top=header_h + answer_h, width=btn_w, height=btn_h,
        font_size=20, font_color=WHITE, bg_color=DARK_BLUE, bold=True,
        border_color=WHITE,
    )

    return slide, board_btn


# ---------------------------------------------------------------------------
# Main presenter
# ---------------------------------------------------------------------------

def create_jeopardy_presentation(categories, values, output_file):
    """
    Build the complete Jeopardy presentation and save to *output_file*.

    Slide layout
    ~~~~~~~~~~~~
    Slide 1         – Main board
    Slides 2, 3     – Clue & Answer for (category 0, value 0)
    Slides 4, 5     – Clue & Answer for (category 0, value 1)
    …
    Slides 60, 61   – Clue & Answer for (category 5, value 4)
    """
    prs = Presentation()

    # 16:9 widescreen (13.33 × 7.5 inches)
    SLIDE_WIDTH_IN = 13.33
    SLIDE_HEIGHT_IN = 7.5
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    # 1. Board slide (shapes collected for hyperlinks later)
    board_slide, value_shapes = build_board_slide(prs, categories, values)

    # 2. Clue + Answer slides for every cell
    #    Iterate category-major so the slide order is predictable.
    pairs = []  # [(clue_slide, answer_btn, answer_slide, board_btn), ...]
    for cat_idx, category in enumerate(categories):
        for val_idx, value in enumerate(values):
            clue_slide, answer_btn = build_clue_slide(prs, category, value)
            answer_slide, board_btn = build_answer_slide(prs, category, value)
            pairs.append((clue_slide, answer_btn, answer_slide, board_btn))

    # 3. Wire up hyperlinks
    #    value_shapes[row_idx][col_idx]  <->  pairs[cat_idx * num_values + val_idx]
    num_values = len(values)
    for cat_idx in range(len(categories)):
        for val_idx in range(num_values):
            cell_num = cat_idx * num_values + val_idx
            clue_slide = pairs[cell_num][0]
            # Board cell → Clue slide
            set_slide_hyperlink(value_shapes[val_idx][cat_idx], clue_slide)

    for clue_slide, answer_btn, answer_slide, board_btn in pairs:
        # Clue   → Answer slide
        set_slide_hyperlink(answer_btn, answer_slide)
        # Answer → Board slide
        set_slide_hyperlink(board_btn, board_slide)

    prs.save(output_file)
    num_slides = len(prs.slides)
    print(f"✓ Saved '{output_file}'  ({num_slides} slides)")
    print(f"  Board: slide 1")
    print(f"  Clue/Answer pairs: slides 2–{num_slides}")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description=(
            "Generate a Jeopardy-style PowerPoint presentation (.pptx).\n\n"
            "Navigation flow inside the deck:\n"
            "  Board cell  →  Clue slide  →  Answer slide  →  Board"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "--categories",
        nargs=6,
        metavar="CATEGORY",
        default=DEFAULT_CATEGORIES,
        help=(
            "Six category names separated by spaces. "
            "Wrap multi-word names in quotes, e.g. "
            '"Science" "World History" "Pop Music"'
        ),
    )
    parser.add_argument(
        "--round",
        dest="round_type",
        choices=["single", "double", "triple"],
        default="single",
        help=(
            "Round type — controls default dollar values:\n"
            "  single  → $200 $400 $600 $800 $1000  (default)\n"
            "  double  → $400 $800 $1200 $1600 $2000\n"
            "  triple  → $600 $1200 $1800 $2400 $3000"
        ),
    )
    parser.add_argument(
        "--values",
        nargs=5,
        type=int,
        metavar="VALUE",
        help=(
            "Five custom dollar values (overrides --round). "
            "Example: --values 100 200 300 400 500"
        ),
    )
    parser.add_argument(
        "--output",
        default="jeopardy.pptx",
        metavar="FILE",
        help="Output filename (default: jeopardy.pptx)",
    )

    args = parser.parse_args()

    categories = args.categories
    values = args.values if args.values else ROUND_VALUES[args.round_type]

    print("Jeopardy PowerPoint Generator")
    print(f"  Categories : {', '.join(categories)}")
    print(f"  Values     : {', '.join(f'${v:,}' for v in values)}")
    print(f"  Output     : {args.output}")
    print()

    create_jeopardy_presentation(categories, values, args.output)


if __name__ == "__main__":
    main()
